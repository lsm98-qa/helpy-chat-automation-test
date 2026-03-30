from pathlib import Path
from zipfile import ZipFile
import xml.etree.ElementTree as ET
import argparse

from pptx import Presentation

# 사용 방법:
# python tests/tools/ppt_create/confirm_ppt_file.py --file <pptx_파일명>

NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def get_pptx_file_by_name(downloads_dir: Path, file_name: str) -> Path:
    pptx_file = downloads_dir / file_name
    if not pptx_file.exists() or not pptx_file.is_file():
        raise RuntimeError(f"파일이 없습니다: {pptx_file}")
    if pptx_file.suffix.lower() != ".pptx":
        raise RuntimeError(f".pptx 파일이 아닙니다: {pptx_file}")
    return pptx_file


def count_slide_count(pptx_path: Path) -> int:
    with ZipFile(pptx_path) as zf:
        presentation_xml = zf.read("ppt/presentation.xml")

    root = ET.fromstring(presentation_xml)

    slide_count = len(root.findall("./p:sldIdLst/p:sldId", NS))
    return slide_count


def count_sections_by_layout(presentation: Presentation) -> int:
    return sum(
        1
        for slide in presentation.slides
        if (slide.slide_layout.name or "").strip().lower() == "section_title"
    )


def iter_slide_titles_and_layouts(presentation: Presentation):
    for idx, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title_text = title_shape.text.strip() if title_shape and title_shape.text else ""
        if not title_text:
            continue
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        yield idx, title_text, layout_name


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--file", required=True, help="검증할 pptx 파일명")
    args = parser.parse_args()

    downloads_dir = Path(__file__).resolve().parent / "downloads"
    pptx_file = get_pptx_file_by_name(downloads_dir, args.file)
    total_slides = count_slide_count(pptx_file)
    presentation = Presentation(str(pptx_file))
    sections = count_sections_by_layout(presentation)
    print(f"total_slides={total_slides}")
    print(f"sections={sections}")
    for slide_idx, title, layout in iter_slide_titles_and_layouts(presentation):
        print(f"slide[{slide_idx}] title={title}")
        print(f"slide[{slide_idx}] layout={layout}")

